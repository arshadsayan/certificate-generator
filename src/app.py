from flask import Flask, request, send_file
from flask_cors import CORS
from pptx import Presentation
import pandas as pd
import io
import os
import zipfile
import datetime
import shutil

app = Flask(__name__)
CORS(app)  # Enable CORS to allow requests from React frontend

# Path to save the certificates
save_path = 'certificates/'

@app.route('/upload', methods=['POST'])
def upload_files():
    try:
        if not os.path.exists(save_path):
            os.makedirs(save_path)

        # Get the files from the request
        ppt_file = request.files['ppt']
        excel_file = request.files['excel']

        # Load the PowerPoint template
        ppt_template = ppt_file.read()

        # Load the Excel file
        sheet = pd.read_excel(io.BytesIO(excel_file.read()))

        # Iterate through each row in Excel
        for index, row in sheet.iterrows():
            stud_name = row.iloc[0]
            course_name = row.iloc[1]
            from_date = row.iloc[2]
            to_date = row.iloc[3]
            org_dept_name = row.iloc[4]
            project_name = row.iloc[5]

            # Load the template for each row
            prs = Presentation(io.BytesIO(ppt_template))
            slide = prs.slides[0]  # Adjust index based on your template

            # Replace placeholders with data from Excel and make the text bold
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if 'Stud_name' in run.text:
                                run.text = run.text.replace('Stud_name', stud_name)
                                run.font.bold = True
                            if 'course_name' in run.text:
                                run.text = run.text.replace('course_name', course_name)
                                run.font.bold = True
                            if 'from_date' in run.text:
                                run.text = run.text.replace('from_date', str(from_date))
                                run.font.bold = True
                            if 'to_date' in run.text:
                                run.text = run.text.replace('to_date', str(to_date))
                                run.font.bold = True
                            if 'org_dept_name' in run.text:
                                run.text = run.text.replace('org_dept_name', org_dept_name)
                                run.font.bold = True
                            if 'project_name' in run.text:
                                run.text = run.text.replace('project_name', project_name)
                                run.font.bold = True

            # Save the filled-out certificate as a new PowerPoint file
            timestamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
            output_file = f"{save_path}Certificate_{stud_name}_{timestamp}.pptx"
            prs.save(output_file)

        # Send back a zip file with all the certificates
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as z:
            for filename in os.listdir(save_path):
                if filename.endswith('.pptx'):
                    z.write(os.path.join(save_path, filename), filename)
        zip_buffer.seek(0)

        # Clean up the certificates directory after creating the zip
        shutil.rmtree(save_path)

        return send_file(zip_buffer, as_attachment=True, download_name='certificates.zip', mimetype='application/zip')

    except Exception as e:
        return str(e), 500

if __name__ == '__main__':
    app.run(debug=True)
